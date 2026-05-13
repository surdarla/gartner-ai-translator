import os
import re
import logging
import traceback
from abc import ABC, abstractmethod

class DocumentProcessor(ABC):
    def __init__(self, translator):
        self.translator = translator

    @abstractmethod
    def process(self, input_path, output_path, progress_callback=None, test_mode=False):
        pass

class PDFProcessor(DocumentProcessor):
    def process(self, input_path, output_path, progress_callback=None, test_mode=False):
        import fitz
        import pikepdf
        
        logging.info(f"Starting PDF translation: {input_path}")
        if progress_callback: progress_callback(0, 100, "PDF 파일 분석 시작...", log_msg="[시스템] PDF 파일 분석을 시작합니다.")
        
        try:
            doc = fitz.open(input_path)
            total_pages = len(doc)
            
            if test_mode:
                total_pages = min(total_pages, 10)
                logging.info(f"Test mode: Processing first {total_pages} pages.")
                if progress_callback: progress_callback(0, 100, log_msg=f"[디버그] 테스트 모드: 상위 {total_pages} 페이지만 처리합니다.")
            
            # Step 1: Extract Text
            all_blocks = []
            if progress_callback: progress_callback(5, 100, "텍스트 데이터 추출 중...", log_msg="[1/4] 텍스트 레이아웃 분석 및 추출을 시작합니다.")
            
            for page_num in range(total_pages):
                page = doc[page_num]
                blocks = page.get_text("dict")["blocks"]
                page_extracted_count = 0
                for b in blocks:
                    if b["type"] == 0:
                        text = " ".join([l["spans"][0]["text"] for l in b["lines"]]).strip()
                        if text and re.search(self.translator.src_regex, text):
                            all_blocks.append({
                                "page": page_num,
                                "bbox": fitz.Rect(b["bbox"]),
                                "text": text
                            })
                            page_extracted_count += 1
                
                if progress_callback and (page_num + 1) % 5 == 0:
                    progress_callback(5 + (page_num/total_pages)*20, 100, f"추출 중: {page_num+1}/{total_pages}p", log_msg=f"[추출] {page_num+1} 페이지 완료 (블록 {page_extracted_count}개)")

            if not all_blocks:
                if progress_callback: progress_callback(100, 100, "번역할 내용 없음", log_msg="[알림] 문서에서 번역 가능한 텍스트 블록을 찾지 못했습니다.")
                doc.save(output_path)
                return True

            if progress_callback: progress_callback(25, 100, "AI 번역 대기 중...", log_msg=f"[2/4] 총 {len(all_blocks)}개의 텍스트 블록 번역을 시작합니다 (AI 요청 중)")

            # Step 2: Translate
            texts_to_translate = [b["text"] for b in all_blocks]
            
            def t_cb(c, t, log_msg="", **kwargs):
                if progress_callback:
                    p = 25 + (c / max(t, 1)) * 50
                    # Use provided log_msg if exists, otherwise fallback to local logic
                    msg = log_msg if log_msg else (f"[번역] AI 응답 수신 중: {c}/{t} 완료" if c > 0 and (c % 20 == 0 or c == t) else "")
                    progress_callback(p, 100, f"AI 번역 진행: {c}/{t}", log_msg=msg)
            
            translated_texts = self.translator.translate_all_concurrent(texts_to_translate, progress_callback=t_cb)
            translated_texts = [t if t else "번역 실패" for t in translated_texts]
            
            for i, b in enumerate(all_blocks):
                b["translated"] = translated_texts[i]

            # Step 3: Draw Text
            if progress_callback: progress_callback(75, 100, "번역 결과 렌더링 중...", log_msg="[3/4] 번역된 내용을 문서에 다시 그리는 중입니다.")
            
            def find_font(tgt_code):
                from core.config import get_back_dir
                back_dir = get_back_dir()
                google_sans = os.path.join(back_dir, "data", "fonts", "Google Sans", "GoogleSans-Regular.ttf")
                if os.path.exists(google_sans): return google_sans
                return None

            font_path = find_font(self.translator.tgt_lang_code)
            
            for page_num in range(total_pages):
                page = doc[page_num]
                if font_path:
                    try:
                        page.insert_font(fontname="F0", fontfile=font_path)
                        fontname = "F0"
                    except: fontname = "helv"
                else: fontname = "helv"

                page_blocks = [b for b in all_blocks if b["page"] == page_num]
                for b in page_blocks:
                    rect = b["bbox"]
                    page.add_redact_annot(rect, fill=(1,1,1))
                    page.apply_redactions()
                    text_to_insert = b["translated"]
                    try:
                        rc = page.insert_textbox(rect, text_to_insert, fontsize=11, fontname=fontname, align=0, color=(0,0,0))
                        if rc < 0:
                            shrunk = self.translator.shrink_text(text_to_insert, 100)
                            page.insert_textbox(rect, shrunk, fontsize=9, fontname=fontname, align=0, color=(0,0,0))
                    except: pass
                
                if progress_callback and (page_num + 1) % 10 == 0:
                    progress_callback(75 + (page_num/total_pages)*20, 100, log_msg=f"[렌더링] {page_num+1} 페이지 완료")

            if test_mode and len(doc) > total_pages:
                doc.delete_pages(total_pages, len(doc) - 1)

            temp_path = output_path + ".tmp.pdf"
            doc.save(temp_path)
            doc.close()

            # Step 4: Stitch Links
            if progress_callback: progress_callback(95, 100, "파일 최종 저장 중...", log_msg="[4/4] 링크 복원 및 파일 최종 저장 단계입니다.")
            orig_pdf = pikepdf.Pdf.open(input_path)
            trans_pdf = pikepdf.Pdf.open(temp_path)
            
            for i, orig_page in enumerate(orig_pdf.pages):
                if i < len(trans_pdf.pages):
                    trans_page = trans_pdf.pages[i]
                    if "/Annots" in orig_page:
                        link_annots = pikepdf.Array()
                        for annot in orig_page.Annots:
                            if annot.get("/Subtype") == "/Link":
                                link_annots.append(trans_pdf.copy_foreign(annot))
                        if len(link_annots) > 0:
                            if "/Annots" not in trans_page: trans_page.Annots = pikepdf.Array()
                            trans_page.Annots.extend(link_annots)
            
            if "/Names" in orig_pdf.Root:
                trans_pdf.Root.Names = trans_pdf.copy_foreign(orig_pdf.Root.Names)

            trans_pdf.save(output_path)
            os.remove(temp_path)
            if progress_callback: progress_callback(100, 100, "번역 완료", log_msg="[성공] 모든 작업이 완료되었습니다. 파일을 다운로드하세요.")
            return True
            
        except Exception as e:
            err_msg = f"에러 발생: {str(e)}\n{traceback.format_exc()}"
            logging.error(err_msg)
            if progress_callback: progress_callback(0, 100, "오류 발생", log_msg=f"[치명적 에러] {str(e)}")
            return False

class PPTXProcessor(DocumentProcessor):
    def process(self, input_path, output_path, progress_callback=None, test_mode=False):
        from pptx import Presentation
        logging.info(f"Starting PPTX translation: {input_path}")
        if progress_callback: progress_callback(0, 100, "PPTX 분석 시작", log_msg="[시스템] PPTX 슬라이드 분석을 시작합니다.")
        
        try:
            prs = Presentation(input_path)
            total_slides = len(prs.slides)
            process_limit = total_slides if not test_mode else min(total_slides, 10)
            
            def get_all_text_frames(shape):
                frames = []
                if getattr(shape, "has_text_frame", False): frames.append(shape.text_frame)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if getattr(cell, "text_frame", None): frames.append(cell.text_frame)
                if hasattr(shape, "shapes"):
                    for child in shape.shapes: frames.extend(get_all_text_frames(child))
                return frames

            all_paragraphs = []
            for i, slide in enumerate(prs.slides):
                if i >= process_limit: break
                for shape in slide.shapes:
                    frames = get_all_text_frames(shape)
                    for frame in frames:
                        for p in frame.paragraphs:
                            text = p.text.strip()
                            if text and re.search(self.translator.src_regex, text):
                                all_paragraphs.append(p)
                if progress_callback and (i + 1) % 5 == 0:
                    progress_callback((i/process_limit)*20, 100, log_msg=f"[추출] {i+1}번 슬라이드 분석 완료")

            if not all_paragraphs:
                if progress_callback: progress_callback(100, 100, log_msg="[알림] 번역할 텍스트가 없습니다.")
                prs.save(output_path)
                return True

            if progress_callback: progress_callback(25, 100, "AI 번역 중...", log_msg=f"[2/2] 총 {len(all_paragraphs)}개 텍스트 AI 번역 시작")

            texts_to_translate = [p.text for p in all_paragraphs]
            def t_cb(c, t, log_msg="", **kwargs):
                if progress_callback:
                    p = 25 + (c / max(t, 1)) * 70
                    # Use provided log_msg or default
                    msg = log_msg if log_msg else f"[번역] AI 처리 중: {c}/{t}"
                    progress_callback(p, 100, log_msg=msg)
                    
            translated_texts = self.translator.translate_all_concurrent(texts_to_translate, progress_callback=t_cb)
            
            for i, p in enumerate(all_paragraphs):
                t = translated_texts[i]
                if t:
                    if len(p.runs) > 0:
                        p.runs[0].text = t
                        for r in p.runs[1:]: r.text = ""
                    else: p.text = t

            prs.save(output_path)
            if progress_callback: progress_callback(100, 100, "완료", log_msg="[성공] PPTX 번역 완료")
            return True
        except Exception as e:
            if progress_callback: progress_callback(0, 100, "오류", log_msg=f"[에러] {str(e)}")
            return False
