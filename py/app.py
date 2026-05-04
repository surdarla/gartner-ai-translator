import streamlit as st
import os
import json
import tempfile
import re
import time
import math
import logging
import pandas as pd
from collections import defaultdict
from pptx import Presentation
from google import genai
from google.genai import types
from dotenv import load_dotenv
from deep_translator import GoogleTranslator
import fitz  # PyMuPDF
import pikepdf
import anthropic
import sqlite3
from streamlit_oauth import OAuth2Component

# SQLite DB Setup
DB_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "logs", "usage.db")
os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
def init_db():
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute("CREATE TABLE IF NOT EXISTS usage_logs (id INTEGER PRIMARY KEY AUTOINCREMENT, timestamp TEXT, username TEXT, action TEXT, provider TEXT, target_lang TEXT, file_type TEXT)")
        conn.commit()
init_db()

def log_usage(username, action, provider="", target_lang="", file_type=""):
    try:
        with sqlite3.connect(DB_PATH) as conn:
            conn.execute("INSERT INTO usage_logs (timestamp, username, action, provider, target_lang, file_type) VALUES (datetime('now','localtime'), ?, ?, ?, ?, ?)", (username, action, provider, target_lang, file_type))
            conn.commit()
    except Exception as e:
        logging.error(f"DB Logging failed: {e}")

# 로깅 기본 설정 (루트 경로에 app.log 누적 저장)
logging.basicConfig(filename=os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "logs", "app.log"), level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s", datefmt="%Y-%m-%d %H:%M:%S")

# 페이지 전역 설정
st.set_page_config(page_title="AI PPT 요약 번역기", page_icon="📈", layout="wide")

# ==============================================================================
# 1. 인증 및 세션 (RBAC) 로직
# ==============================================================================
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
if "role" not in st.session_state:
    st.session_state.role = None
if "username" not in st.session_state:
    st.session_state.username = None


def do_login(user, pw):
    if user == "user" and pw == "1234":
        st.session_state.logged_in = True
        st.session_state.role = "user"
        st.session_state.username = user
        logging.info(f"사용자 로그인 성공: {user}")
        return True
    elif user == "admin" and pw == "admin123":
        st.session_state.logged_in = True
        st.session_state.role = "admin"
        st.session_state.username = user
        logging.info(f"관리자 로그인 성공: {user}")
        return True
    return False


def do_logout():
    logging.info(f"로그아웃: {st.session_state.username}")
    st.session_state.logged_in = False
    st.session_state.role = None
    st.session_state.username = None


# 로그인 가드 (미로그인 시 차단)
if not st.session_state.logged_in:
    st.title("🔒 로그인 (Gartner AI Translator)")
    
    tab1, tab2, tab3 = st.tabs(["일반 로그인", "Google SSO", "Microsoft SSO"])
    
    with tab1:
        with st.form("login_form"):
            st.markdown("사내 번역 시스템에 접근하려면 소속 및 권한 확인을 위해 로그인하세요.")
            st.info("💡 접속 가이드 (체험판)\n- 👨‍💻 일반 사용자: ID `user` / PW `1234`\n- 👑 시스템 관리자: ID `admin` / PW `admin123`")
            u_input = st.text_input("Username")
            p_input = st.text_input("Password", type="password")
            submit = st.form_submit_button("로그인")
            if submit:
                if do_login(u_input, p_input):
                    st.rerun()
                else:
                    logging.warning(f"로그인 실패 시도: {u_input}")
                    st.error("❌ 아이디 또는 비밀번호가 올바르지 않습니다.")
                    
    with tab2:
        st.subheader("Google 계정으로 로그인")
        GOOGLE_CLIENT_ID = os.environ.get("GOOGLE_CLIENT_ID", "")
        GOOGLE_CLIENT_SECRET = os.environ.get("GOOGLE_CLIENT_SECRET", "")
        if GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET:
            google_oauth2 = OAuth2Component(GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET, "https://accounts.google.com/o/oauth2/v2/auth", "https://oauth2.googleapis.com/token", "https://oauth2.googleapis.com/token", "https://oauth2.googleapis.com/revoke")
            result = google_oauth2.authorize_button("Google Login", "http://localhost:8501/oauth2callback", scope="email profile")
            if result and "token" in result:
                st.session_state.logged_in = True
                st.session_state.role = "user" # Default role for SSO
                st.session_state.username = "Google User" # In a real app, decode JWT to get email
                st.rerun()
        else:
            st.warning("⚠️ Google SSO가 설정되지 않았습니다. 관리자에게 문의하세요 (.env 파일에 GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET 설정 필요).")

    with tab3:
        st.subheader("Microsoft 계정으로 로그인")
        MS_CLIENT_ID = os.environ.get("MICROSOFT_CLIENT_ID", "")
        MS_CLIENT_SECRET = os.environ.get("MICROSOFT_CLIENT_SECRET", "")
        MS_TENANT_ID = os.environ.get("MICROSOFT_TENANT_ID", "common")
        if MS_CLIENT_ID and MS_CLIENT_SECRET:
            ms_oauth2 = OAuth2Component(MS_CLIENT_ID, MS_CLIENT_SECRET, f"https://login.microsoftonline.com/{MS_TENANT_ID}/oauth2/v2.0/authorize", f"https://login.microsoftonline.com/{MS_TENANT_ID}/oauth2/v2.0/token", f"https://login.microsoftonline.com/{MS_TENANT_ID}/oauth2/v2.0/token", "")
            result = ms_oauth2.authorize_button("Microsoft Login", "http://localhost:8501/oauth2callback", scope="openid email profile")
            if result and "token" in result:
                st.session_state.logged_in = True
                st.session_state.role = "user" # Default role for SSO
                st.session_state.username = "Microsoft User" # In a real app, decode JWT to get email
                st.rerun()
        else:
            st.warning("⚠️ Microsoft SSO가 설정되지 않았습니다. 관리자에게 문의하세요 (.env 파일에 MICROSOFT_CLIENT_ID, MICROSOFT_CLIENT_SECRET 설정 필요).")
    st.stop()


# ==============================================================================
# 2. 유틸리티 함수 모음
# ==============================================================================
def load_glossary():
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "edtech_glossary.json")
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def get_all_text_frames(shape):
    frames = []
    if getattr(shape, "has_text_frame", False):
        frames.append(shape.text_frame)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if getattr(cell, "text_frame", None):
                    frames.append(cell.text_frame)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            frames.extend(get_all_text_frames(child))
    return frames


def translate_batch_gemini(client, batch_texts, system_instruction, src_lang_name="Korean", lang_name="Business English"):
    if not batch_texts:
        return batch_texts

    input_dict = {str(i): text for i, text in enumerate(batch_texts)}

    prompt = f"Translate this JSON object from {src_lang_name} into {lang_name}. You MUST output ONLY a valid JSON object where the EXACT SAME keys are mapped to their translated/summarized versions. Do not skip any keys.\n\n"
    prompt += json.dumps(input_dict, ensure_ascii=False)

    try:
        response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt, config=types.GenerateContentConfig(system_instruction=system_instruction, temperature=0.0, response_mime_type="application/json"))
        response_json = json.loads(response.text)

        translated_list = []
        for i in range(len(batch_texts)):
            translated_list.append(response_json.get(str(i), batch_texts[i]))

        return translated_list

    except Exception as e:
        print(f"  --> 일괄 번역 실패(원본 유지): {e}")
        return batch_texts

def translate_batch_claude(client, batch_texts, system_instruction, src_lang_name="Korean", lang_name="Business English"):
    if not batch_texts:
        return batch_texts

    input_dict = {str(i): text for i, text in enumerate(batch_texts)}

    prompt = f"Translate this JSON object from {src_lang_name} into {lang_name}. You MUST output ONLY a valid JSON object where the EXACT SAME keys are mapped to their translated/summarized versions. Do not skip any keys.\n\n"
    prompt += json.dumps(input_dict, ensure_ascii=False)

    try:
        response = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=4096,
            temperature=0.0,
            system=system_instruction,
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        response_text = response.content[0].text
        # Claude might wrap JSON in markdown blocks
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()
        
        response_json = json.loads(response_text)

        translated_list = []
        for i in range(len(batch_texts)):
            translated_list.append(response_json.get(str(i), batch_texts[i]))

        return translated_list

    except Exception as e:
        print(f"  --> 일괄 번역 실패(원본 유지): {e}")
        return batch_texts

def translate_batch_free(batch_texts, src_code, tgt_code, glossary):
    """src_code/tgt_code: ISO 639-1 코드 (ko, en, ja, auto)"""
    translator = GoogleTranslator(source=src_code, target=tgt_code)

    translated_texts = []
    for text in batch_texts:
        # 사전 치환 (고유명사 오역 방지용 대괄호 마킹)
        temp_text = text
        for k, v in glossary.items():
            temp_text = temp_text.replace(k, f"[{v}]")

        try:
            res = translator.translate(temp_text)
            if res:
                res = res.replace("[", "").replace("]", "")
                translated_texts.append(res)
            else:
                translated_texts.append(text)
        except Exception as e:
            logging.error(f"무료 번역 에러: {e}")
            translated_texts.append(text)

    return translated_texts


# ==============================================================================
# PDF 번역 전용 유틸리티 함수
# ==============================================================================

def stitch_links_with_pikepdf(orig_path, trans_path, out_path):
    """
    pikepdf를 사용하여 원본 PDF의 모든 어노테이션(링크)과 목적지(Names/Dests)를 
    번역된 PDF로 그대로 복사(Stitching)합니다.
    """
    try:
        with pikepdf.Pdf.open(orig_path) as orig_pdf, pikepdf.Pdf.open(trans_path) as trans_pdf:
            for i in range(min(len(orig_pdf.pages), len(trans_pdf.pages))):
                orig_page = orig_pdf.pages[i]
                trans_page = trans_pdf.pages[i]
                if hasattr(orig_page, 'Annots'):
                    try:
                        trans_page.Annots = trans_pdf.copy_foreign(orig_page.Annots)
                    except:
                        if '/Annots' in orig_page.obj:
                            if '/Annots' in trans_page.obj: del trans_page.obj['/Annots']
                            trans_page.obj['/Annots'] = trans_pdf.copy_foreign(orig_page.obj['/Annots'])
            
            orig_root = orig_pdf.trailer.Root
            trans_root = trans_pdf.trailer.Root
            
            for key in ['/Names', '/Dests', '/ViewerPreferences']:
                if key in orig_root:
                    try:
                        if key in trans_root: del trans_root[key]
                        trans_root[key] = trans_pdf.copy_foreign(orig_root[key])
                    except Exception as e:
                        if "direct object handle" in str(e).lower():
                            try:
                                orig_val = orig_root[key]
                                if isinstance(orig_val, pikepdf.Dictionary):
                                    new_dict = pikepdf.Dictionary()
                                    for subkey, subval in orig_val.items():
                                        try: new_dict[subkey] = trans_pdf.copy_foreign(subval)
                                        except: pass
                                    trans_root[key] = new_dict
                            except: pass
                        else:
                            logging.warning(f"Error copying global key {key}: {e}")
            
            trans_pdf.save(out_path)
            return True
    except Exception as e:
        logging.error(f"Error during link stitching in app: {e}")
        return False


def find_google_sans_font(tgt_code="en"):
    """시스템에서 적절한 폰트 탐색 (사용자 요청에 따라 Google Sans 최우선, 단 CJK는 검증 후 Fallback)"""
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    google_sans_path = os.path.join(base_dir, "font", "Google Sans", "GoogleSans-Regular.ttf")
    
    is_cjk = tgt_code in ["ko", "ja"]
    
    if os.path.exists(google_sans_path):
        if is_cjk and os.path.getsize(google_sans_path) < 1024 * 1024:
            logging.debug(f"로컬 Google Sans가 너무 작음({os.path.getsize(google_sans_path)} bytes). CJK Fallback 수행.")
            pass
        else:
            return google_sans_path

    if tgt_code == "ko":
        candidates = [
            "/System/Library/Fonts/AppleSDGothicNeo.ttc",
            "/Library/Fonts/NanumGothic.ttf",
            "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
        ]
    elif tgt_code == "ja":
        candidates = [
            "/System/Library/Fonts/Hiragino Sans GB.ttc",
            "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
        ]
    else:
        candidates = [
            os.path.expanduser("~/Library/Fonts/GoogleSans-Regular.ttf"),
            "/Library/Fonts/GoogleSans-Regular.ttf",
        ]
    
    for p in candidates:
        if os.path.exists(p):
            return p
    
    fallback_all = [
        "/System/Library/Fonts/AppleSDGothicNeo.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/STHeiti Light.ttc",
    ]
    for p in fallback_all:
        if os.path.exists(p):
            return p
            
    return None


def shrink_text_gemini(client, text, target_len, lang_name):
    """번역 결과가 너무 길 때 핵심 보존 재압축 API 호출"""
    prompt = (
        f"The following {lang_name} text is too long for a fixed-size PDF text box.\n"
        f"Compress it to strictly within {target_len} characters while keeping the core meaning.\n"
        f"Output ONLY the final compressed text, no explanation.\n\nText: {text}"
    )
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(temperature=0.0),
        )
        return response.text.strip()
    except Exception as e:
        logging.error(f"Shrink 압축 실패: {e}")
        return text

def shrink_text_claude(client, text, target_len, lang_name):
    prompt = (
        f"The following {lang_name} text is too long for a fixed-size PDF text box.\n"
        f"Compress it to strictly within {target_len} characters while keeping the core meaning.\n"
        f"Output ONLY the final compressed text, no explanation.\n\nText: {text}"
    )
    try:
        response = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=1024,
            temperature=0.0,
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        return response.content[0].text.strip()
    except Exception as e:
        logging.error(f"Shrink 압축 실패: {e}")
        return text


def translate_pdf(tmp_in_path, output_path, client, system_instruction,
                  src_code, tgt_code, src_regex, glossary, llm_provider,
                  my_bar, username, src_lang_name, lang_name):
    """PDF 번역 메인 함수 (Redact → Google Sans 삽입 + 2-Pass Shrink Loop)"""
    doc = fitz.open(tmp_in_path)
    font_path = find_google_sans_font(tgt_code)
    FONT_NAME = "f0"
    SHRINK_THRESHOLD = 1.35

    # 1단계: 전체 텍스트 블록 수집
    all_blocks = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        links = page.get_links()
        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b.get("type") != 0: continue
            current_p = {"text": "", "bbox": None, "fontsize": None, "color": None, "is_link": False}
            for line in b["lines"]:
                for span in line["spans"]:
                    span_text = span["text"].strip()
                    if not span_text: continue
                    span_rect = fitz.Rect(span["bbox"])
                    line_fs = span["size"]
                    line_color = span["color"]
                    
                    is_link_span = any(fitz.Rect(l["from"]).intersects(span_rect) for l in links if l.get("kind") in (1, 2, 4))
                    is_list_item = bool(re.match(r"^\s*(?:Step\s*\d+:?|\d+[\.\)]|\-|•|\*|[A-Za-z]\.|[①-⑳]|[⑴-⑽])\s+", span_text))
                    
                    if current_p["fontsize"] is None:
                        current_p = {"text": span["text"], "bbox": span_rect, "fontsize": line_fs, "color": line_color, "is_link": is_link_span}
                    else:
                        color_differs = current_p["color"] != line_color
                        fs_differs = abs(current_p["fontsize"] - line_fs) > 1.5
                        link_differs = current_p["is_link"] != is_link_span
                        horizontal_gap = span_rect.x0 - current_p["bbox"].x1
                        vertical_gap = span_rect.y0 - current_p["bbox"].y1
                        
                        is_far_horizontal = (horizontal_gap > line_fs * 1.5) and (abs(span_rect.y0 - current_p["bbox"].y0) < 5)
                        ends_with_period = current_p["text"].strip().endswith(('。', '.', '！', '？', ':'))
                        is_far_vertical = (vertical_gap > line_fs * 0.8) or (vertical_gap > 0 and ends_dot)
                        
                        if (color_differs or fs_differs or link_differs or is_list_item or is_far_horizontal or is_far_vertical):
                            if current_p["text"].strip() and re.search(src_regex, current_p["text"]):
                                all_blocks.append({
                                    "page": page_num, "bbox": current_p["bbox"], "text": current_p["text"].strip(),
                                    "fontsize": current_p["fontsize"], "color": current_p["color"], "is_link": current_p["is_link"]
                                })
                            current_p = {"text": span["text"], "bbox": span_rect, "fontsize": line_fs, "color": line_color, "is_link": is_link_span}
                        else:
                            curr_text = current_p["text"]
                            if curr_text.endswith("-"):
                                current_p["text"] = curr_text[:-1] + span["text"]
                            else:
                                if curr_text and not curr_text.endswith(" ") and not span["text"].startswith(" "):
                                    current_p["text"] += " "
                                current_p["text"] += span["text"]
                            current_p["bbox"] |= span_rect
            if current_p["fontsize"] is not None and current_p["text"].strip() and re.search(src_regex, current_p["text"]):
                all_blocks.append({
                    "page": page_num, "bbox": current_p["bbox"], "text": current_p["text"].strip(),
                    "fontsize": current_p["fontsize"], "color": current_p["color"], "is_link": current_p["is_link"]
                })

    total = len(all_blocks)
    if total == 0: return 0

    # 2단계: 배치 번역
    batch_size = 50
    completed = 0
    translated_results = []
    for i in range(0, total, batch_size):
        batch = all_blocks[i : i + batch_size]
        batch_texts = [b["text"] for b in batch]
        if llm_provider == "Free (Google Translator)":
            translated = translate_batch_free(batch_texts, src_code, tgt_code, glossary)
        elif llm_provider == "Claude":
            translated = translate_batch_claude(client, batch_texts, system_instruction, src_lang_name, lang_name)
        else:
            translated = translate_batch_gemini(client, batch_texts, system_instruction, src_lang_name, lang_name)
        translated_results.extend(translated)
        completed += len(batch)
        my_bar.progress(min(completed / total * 0.65, 0.65), text=f"⏳ 1차 번역 중... [{completed}/{total}] 단락")
        if not use_free_model and i + batch_size < total: time.sleep(5)

    # 3단계: 2차 압축(Shrink Loop)
    shrink_count = 0
    for idx, (block_info, translated) in enumerate(zip(all_blocks, translated_results)):
        original_len = max(len(block_info["text"]), 1)
        if llm_provider != "Free (Google Translator)" and len(translated) > original_len * SHRINK_THRESHOLD:
            target_len = int(original_len * 1.15)
            if llm_provider == "Claude":
                compressed = shrink_text_claude(client, translated, target_len, lang_name)
            else:
                compressed = shrink_text_gemini(client, translated, target_len, lang_name)
            translated_results[idx] = compressed
            shrink_count += 1
        my_bar.progress(0.65 + min((idx + 1) / total * 0.2, 0.2), text=f"🔬 2차 길이 검사 중... [{idx+1}/{total}] (압축 발동: {shrink_count}회)")

    # 4단계: 페이지별 Redact → 번역 텍스트 삽입
    page_group = defaultdict(list)
    for block_info, translated in zip(all_blocks, translated_results):
        page_group[block_info["page"]].append((block_info, translated))

    processed_pages = 0
    for page_num, items in page_group.items():
        page = doc[page_num]
        for block_info, _ in items:
            page.add_redact_annot(block_info["bbox"], fill=None)
        page.apply_redactions(images=0, graphics=0, text=0)

        current_font = "helv"
        if font_path:
            try:
                page.insert_font(fontname=FONT_NAME, fontfile=font_path)
                current_font = FONT_NAME
            except: pass

        for block_info, translated in items:
            c = block_info["color"]
            color = ((c >> 16 & 0xFF) / 255, (c >> 8 & 0xFF) / 255, (c & 0xFF) / 255)
            try:
                rect = block_info["bbox"]
                original_fs = block_info["fontsize"]
                min_fs = 5
                inserted = False
                for fs in range(int(math.ceil(original_fs)), min_fs - 1, -1):
                    shape = page.new_shape()
                    rc = shape.insert_textbox(rect, translated, fontsize=fs, fontname=current_font, color=color, align=0)
                    if rc >= 0:
                        shape.commit()
                        inserted = True
                        break
                if not inserted:
                    shape = page.new_shape()
                    shape.insert_textbox(rect, translated, fontsize=min_fs, fontname=current_font, color=color, align=0)
                    shape.commit()
            except: pass
        processed_pages += 1
        my_bar.progress(0.85 + min(processed_pages / len(page_group) * 0.15, 0.15), text=f"🎨 텍스트 렌더링 중... [{processed_pages}/{len(page_group)}] 페이지")

    render_temp = output_path + ".render_tmp.pdf"
    doc.save(render_temp, garbage=4, deflate=True)
    doc.close()
    
    if stitch_links_with_pikepdf(tmp_in_path, render_temp, output_path):
        if os.path.exists(render_temp): os.remove(render_temp)
        logging.info(f"✨ 번역 및 링크 복제 완료: {output_path}")
    else:
        os.rename(render_temp, output_path)
        logging.warning("⚠️ 링크 복제 실패")

    return total


# ==============================================================================
# 3. 메인 서비스 애플리케이션 화면 구성
# ==============================================================================

# 3-1. 사이드바 구성 (RBAC 표시 및 옵션)
st.sidebar.markdown(f"👤 접속자: `{st.session_state.username}`\n\n👑 권한: `{st.session_state.role.upper()}`")
if st.sidebar.button("🚪 로그아웃", width="stretch"):
    do_logout()
    st.rerun()
st.sidebar.markdown("---")

st.sidebar.header("⚙️ 환경 설정")
load_dotenv()
llm_provider = st.sidebar.selectbox("🤖 번역 AI 선택", ["Free (Google Translator)", "Gemini", "Claude"])
api_key = ""
if llm_provider != "Free (Google Translator)":
    api_key = st.sidebar.text_input(f"🔑 {llm_provider} API Key 지정", value="", type="password", help=f"입력하지 않으면 번역을 시작할 수 없습니다.")

st.sidebar.markdown("---")
target_language = st.sidebar.selectbox(
    "🌐 번역 방향 선택",
    options=[
        "한국어 → English",
        "한국어 → 日本語",
        "English → 한국어",
        "日本語 → 한국어",
    ]
)


def render_translation_view():
    st.title("📈 사내 전용 AI PPT / PDF 자동 요약 번역기")
    st.markdown(
        """
    파일을 드래그 앤 드롭하면, AI가 주변 문맥을 분석하여 원본 공간의 크기에 맞게 핵심만 뽑아 요약 번역을 진행합니다.
    - 글자 넘침(Overflow) 방지: 외국어 번역 시 길이가 팽창하여 도형 밖으로 나가는 현상을 막습니다.
    - 글자 크기 유지 보장: 글꼴 크기를 억지로 줄이지 않고, 문장을 똑똑하게 압축(Summarize)합니다.
    - PDF 2-Pass 압축: PDF는 바운딩박스 크기를 직접 측정하여, 넘칠 경우 AI가 핵심만 남겨 재압축합니다.
    """
    )

    uploaded_file = st.file_uploader("📥 여기에 사내 PPTX 또는 PDF 파일을 드래그 하세요", type=["pptx", "pdf"])

    if uploaded_file is not None:
        file_ext = uploaded_file.name.lower().rsplit(".", 1)[-1]
        is_pdf = file_ext == "pdf"

        btn_label = "🚀 문맥 기반 AI 요약 번역 시작 (PDF 2-Pass 압축 포함)" if is_pdf else "🚀 문맥 기반 AI 요약 번역 시작"
        if st.button(btn_label, width="stretch"):
            if llm_provider != "Free (Google Translator)" and not api_key:
                st.error(f"❌ {llm_provider} API Key를 입력해주세요.")
                st.stop()

            logging.info(f"[{st.session_state.username}] 번역 시작 (파일: {uploaded_file.name}, 언어: {target_language}, 형식: {file_ext.upper()}, AI: {llm_provider})")
            log_usage(st.session_state.username, "Translate Start", llm_provider, target_language, file_ext.upper())
            my_bar = st.progress(0, text="파일 분석 및 텍스트 추출 중...")

            suffix = f".{file_ext}"
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_in:
                tmp_in.write(uploaded_file.getvalue())
                tmp_in_path = tmp_in.name

            output_path = tmp_in_path.replace(suffix, f"_Translated{suffix}")

            client = None
            if llm_provider == "Gemini":
                client = genai.Client(api_key=api_key)
            elif llm_provider == "Claude":
                client = anthropic.Anthropic(api_key=api_key)

            glossary = load_glossary()
            glossary_text = "\n".join([f"- {k}: {v}" for k, v in glossary.items()])

            direction_map = {
                "한국어 → English":  {"src_code": "ko", "tgt_code": "en", "src_lang_name": "Korean",   "lang_name": "Business English",  "src_regex": "[가-힣]"},
                "한국어 → 日本語":  {"src_code": "ko", "tgt_code": "ja", "src_lang_name": "Korean",   "lang_name": "Business Japanese", "src_regex": "[가-힣]"},
                "English → 한국어": {"src_code": "en", "tgt_code": "ko", "src_lang_name": "English",  "lang_name": "Professional Korean",  "src_regex": "[a-zA-Z]"},
                "日本語 → 한국어": {"src_code": "ja", "tgt_code": "ko", "src_lang_name": "Japanese", "lang_name": "Professional Korean",  "src_regex": "[ぁ-ゖァ-ヶ一-鿿]"},
            }
            d = direction_map.get(target_language, direction_map["한국어 → English"])
            src_code     = d["src_code"]
            tgt_code     = d["tgt_code"]
            src_lang_name = d["src_lang_name"]
            lang_name    = d["lang_name"]
            src_regex    = d["src_regex"]

            to_english = tgt_code == "en"
            to_japanese = tgt_code == "ja"
            file_type_label = "PDF text box" if is_pdf else "PPT text box"

            if to_english:
                expansion_logic = f"{src_lang_name} to English translation tends to EXPAND length."
                style_logic = "Omit filler words, concise bullet-points, business English."
                terminology_note = ""
            elif to_japanese:
                expansion_logic = "Use 'Taigen-dome' (体言止め)."
                style_logic = "Professional Japanese analytical terminology."
                terminology_note = "Katakana or Kanji."
            else:
                expansion_logic = "Match original visual space."
                style_logic = "반드시 가트너 에널리스트 스타일의 전문적 비즈니스 한국어(-다, -이다) 사용."
                terminology_note = "의미 모호 시 영문 병기 (예: SASE)."

            system_instruction = f"""You are a Senior Gartner Analyst.
Translate the provided {src_lang_name} texts into highly professional {lang_name}.
- {expansion_logic}
- {style_logic}
- Apply Terminology: {glossary_text}
"""

            try:
                if is_pdf:
                    total_spans = translate_pdf(
                        tmp_in_path, output_path, client, system_instruction,
                        src_code, tgt_code, src_regex, glossary, llm_provider,
                        my_bar, st.session_state.username, src_lang_name, lang_name
                    )
                    if total_spans == 0:
                        st.warning("번역할 텍스트를 찾지 못했습니다.")
                        st.stop()

                    my_bar.progress(1.0, text="✅ PDF 생성 완료!")
                    st.success(f"🎉 총 {total_spans}개 블록 번역 삽입 완료!")

                    with open(output_path, "rb") as file:
                        st.download_button("📤 번역 완료 PDF 다운로드", data=file, file_name=f"Translated_{uploaded_file.name}", mime="application/pdf")

                else:
                    prs = Presentation(tmp_in_path)
                    paragraphs_to_translate = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            for tf in get_all_text_frames(shape):
                                for p in tf.paragraphs:
                                    original_text = p.text.strip()
                                    if original_text and re.search(src_regex, original_text):
                                        paragraphs_to_translate.append((p, original_text))

                    total_paragraphs = len(paragraphs_to_translate)
                    if total_paragraphs == 0:
                        st.warning("번역할 텍스트를 찾지 못했습니다.")
                        st.stop()

                    batch_size = 50
                    completed_paragraphs = 0
                    for i in range(0, total_paragraphs, batch_size):
                        batch_data = paragraphs_to_translate[i : i + batch_size]
                        batch_texts = [info[1] for info in batch_data]

                        if llm_provider == "Free (Google Translator)":
                            translated_texts = translate_batch_free(batch_texts, src_code, tgt_code, glossary)
                        elif llm_provider == "Claude":
                            translated_texts = translate_batch_claude(client, batch_texts, system_instruction, src_lang_name, lang_name)
                        else:
                            translated_texts = translate_batch_gemini(client, batch_texts, system_instruction, src_lang_name, lang_name)

                        for idx, (p, original) in enumerate(batch_data):
                            new_text = translated_texts[idx]
                            if p.runs:
                                p.runs[0].text = new_text
                                for j in range(len(p.runs) - 1, 0, -1):
                                    r = p.runs[j]._r
                                    r.getparent().remove(r)
                            else:
                                p.text = new_text

                        completed_paragraphs += len(batch_data)
                        prog = min(completed_paragraphs / total_paragraphs, 1.0)
                        my_bar.progress(prog, text=f"⏳ 번역 중... [{completed_paragraphs}/{total_paragraphs}]")
                        if llm_provider != "Free (Google Translator)" and i + batch_size < total_paragraphs:
                            time.sleep(5)

                    prs.save(output_path)
                    my_bar.progress(1.0, text="✅ PPTX 생성 완료!")
                    st.success(f"🎉 {total_paragraphs}개 단락 번역 완료!")

                    with open(output_path, "rb") as file:
                        st.download_button("📤 번역 완료 PPT 다운로드", data=file, file_name=f"Translated_{uploaded_file.name}", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

            except Exception as e:
                logging.error(f"번역 에러: {e}")
                st.error(f"번역 중 오류가 발생했습니다: {e}")


def render_admin_dashboard():
    st.title("📦 중앙 관리자 대시보드")
    
    st.subheader("1. 사용자 이용 통계 (DB Logs)")
    try:
        with sqlite3.connect(DB_PATH) as conn:
            df_usage = pd.read_sql_query("SELECT * FROM usage_logs ORDER BY timestamp DESC LIMIT 100", conn)
            st.dataframe(df_usage, width="stretch", height=200)
    except Exception as e:
        st.warning(f"DB 로그를 불러오는 중 오류가 발생했습니다: {e}")

    st.subheader("2. 애플리케이션 시스템 로그")
    log_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "logs", "app.log")
    if os.path.exists(log_path):
        try:
            with open(log_path, "r", encoding="utf-8") as f:
                logs = f.readlines()

            parsed_logs = []
            for raw_log in logs:
                parts = raw_log.split(" | ", 2)
                if len(parts) == 3:
                    parsed_logs.append({"Time": parts[0].strip(), "Level": parts[1].strip(), "Message": parts[2].strip()})

            if parsed_logs:
                df = pd.DataFrame(parsed_logs)
                df = df.sort_values(by="Time", ascending=False).reset_index(drop=True)
                st.dataframe(df, width="stretch", height=200)
            else:
                st.info("로그가 없습니다.")
        except Exception as e:
            st.error(f"로그 읽기 실패: {e}")
    else:
        st.warning("로그 파일이 없습니다.")

    st.subheader("3. 에듀테크 용어집 (Glossary) 관리")
    glossary_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "edtech_glossary.json")
    if os.path.exists(glossary_path):
        with open(glossary_path, "r", encoding="utf-8") as f:
            glossary_data = json.load(f)
        
        # Display as a dataframe for easier editing
        df_glossary = pd.DataFrame(list(glossary_data.items()), columns=["원본 용어", "번역어"])
        edited_df = st.data_editor(df_glossary, num_rows="dynamic", use_container_width=True)
        
        if st.button("용어집 저장하기"):
            new_glossary = dict(zip(edited_df["원본 용어"], edited_df["번역어"]))
            with open(glossary_path, "w", encoding="utf-8") as f:
                json.dump(new_glossary, f, ensure_ascii=False, indent=4)
            st.success("✅ 용어집이 성공적으로 업데이트되었습니다!")
            log_usage(st.session_state.username, "Update Glossary")
    else:
        st.warning("용어집 파일을 찾을 수 없습니다.")


# 4. 권한(Role)에 따른 렌더링 분기
if st.session_state.role == "admin":
    tab1, tab2 = st.tabs(["🚀 요약 번역 서비스", "📦 시스템 로그 대시보드 (Admin)"])
    with tab1:
        render_translation_view()
    with tab2:
        render_admin_dashboard()
else:
    render_translation_view()
