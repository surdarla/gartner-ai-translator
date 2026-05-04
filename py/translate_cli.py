#!/usr/bin/env python3
"""
translate_cli.py — Streamlit 없이 터미널에서 직접 실행하는 AI 번역 스크립트 (Link Stitching 최적화 버전)

사용법:
  uv run python translate_cli.py <파일경로> [옵션]
"""

import os
import sys
import json
import re
import time
from datetime import datetime
import math
import logging
import argparse
import pikepdf
from collections import defaultdict
from pathlib import Path

from dotenv import load_dotenv
from tqdm import tqdm
from concurrent.futures import ThreadPoolExecutor, as_completed

try:
    from pptx import Presentation
except ImportError:
    print("❌ python-pptx 미설치: uv add python-pptx")
    sys.exit(1)

try:
    import fitz  # PyMuPDF
except ImportError:
    print("❌ pymupdf 미설치: uv add pymupdf")
    sys.exit(1)

try:
    from google import genai
    from google.genai import types
except ImportError:
    print("❌ google-genai 미설치: uv add google-genai")
    sys.exit(1)

try:
    import anthropic
except ImportError:
    print("❌ anthropic 미설치: uv add anthropic")
    sys.exit(1)

try:
    from deep_translator import GoogleTranslator
except ImportError:
    print("❌ deep-translator 미설치: uv add deep-translator")
    sys.exit(1)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
    handlers=[
        logging.FileHandler(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "logs", "app.log"), encoding="utf-8"),
        logging.StreamHandler(sys.stdout),
    ],
)

DIRECTION_MAP = {
    "한국어 → English":  {"src_code": "ko", "tgt_code": "en", "src_lang_name": "Korean",   "lang_name": "Business English",  "src_regex": "[가-힣]"},
    "한국어 → 日本語":   {"src_code": "ko", "tgt_code": "ja", "src_lang_name": "Korean",   "lang_name": "Business Japanese", "src_regex": "[가-힣]"},
    "English → 한국어":  {"src_code": "en", "tgt_code": "ko", "src_lang_name": "English",  "lang_name": "Professional Korean", "src_regex": "[a-zA-Z]"},
    "日本語 → 한국어":   {"src_code": "ja", "tgt_code": "ko", "src_lang_name": "Japanese", "lang_name": "Professional Korean", "src_regex": "[ぁ-ゖァ-ヶ一-鿿]"},
}

SHRINK_THRESHOLD = 1.35
BATCH_SIZE = 100
SLEEP_BETWEEN_BATCHES = 0.1

def load_glossary(path=None):
    if path is None:
        path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "edtech_glossary.json")
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def stitch_links_with_pikepdf(orig_path, trans_path, out_path):
    """pikepdf를 이용한 원본 링크 및 전역 Dests 복사 (100% 보존 최적화)"""
    try:
        # qpdf_object_handle 에러 방지를 위해 trailer.Root를 통해 접근
        with pikepdf.Pdf.open(orig_path) as orig_pdf, pikepdf.Pdf.open(trans_path) as trans_pdf:
            # 1. 페이지별 어노테이션(링크) 복사
            for i in range(min(len(orig_pdf.pages), len(trans_pdf.pages))):
                orig_page = orig_pdf.pages[i]
                trans_page = trans_pdf.pages[i]
                
                # 'Annots' 속성을 직접 사용 (pikepdf 9.0+ 권장)
                if hasattr(orig_page, 'Annots'):
                    try:
                        trans_page.Annots = trans_pdf.copy_foreign(orig_page.Annots)
                    except:
                        # Fallback: Dictionary 키로 접근
                        if '/Annots' in orig_page.obj:
                            if '/Annots' in trans_page.obj: del trans_page.obj['/Annots']
                            trans_page.obj['/Annots'] = trans_pdf.copy_foreign(orig_page.obj['/Annots'])
            
            # 2. 전역 목적지 사전 (Names/Dests) 복사
            orig_root = orig_pdf.trailer.Root
            trans_root = trans_pdf.trailer.Root
            
            # /Names, /Dests 등 전역 사전 복사 시도
            for key in ['/Names', '/Dests', '/ViewerPreferences']:
                if key in orig_root:
                    try:
                        # 기존 항목 제거 후 복사
                        if key in trans_root:
                            del trans_root[key]
                        trans_root[key] = trans_pdf.copy_foreign(orig_root[key])
                    except Exception as e:
                        # 특정 파일에서 direct object handle 에러 발생 시 최후의 수단으로 개별 키 복사 시도
                        if "direct object handle" in str(e).lower():
                            try:
                                # Dictionary인 경우 내부 아이템을 하나씩 복사 (Names 하위 등)
                                orig_val = orig_root[key]
                                if isinstance(orig_val, pikepdf.Dictionary):
                                    new_dict = pikepdf.Dictionary()
                                    for subkey, subval in orig_val.items():
                                        try:
                                            new_dict[subkey] = trans_pdf.copy_foreign(subval)
                                        except: pass
                                    trans_root[key] = new_dict
                            except:
                                logging.warning(f"Failed to copy global key {key} completely.")
                        else:
                            logging.warning(f"Error copying global key {key}: {e}")
            
            trans_pdf.save(out_path)
            return True
    except Exception as e:
        print(f"Error during link stitching: {e}")
        return False

def find_google_sans_font(tgt_code="en"):
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    google_sans_path = os.path.join(base_dir, "font", "Google Sans", "GoogleSans-Regular.ttf")
    is_cjk = tgt_code in ["ko", "ja"]
    if os.path.exists(google_sans_path):
        if is_cjk and os.path.getsize(google_sans_path) < 1024 * 1024:
            pass
        else:
            return google_sans_path
    if tgt_code == "ko":
        candidates = ["/System/Library/Fonts/AppleSDGothicNeo.ttc", "/Library/Fonts/NanumGothic.ttf"]
    elif tgt_code == "ja":
        candidates = ["/System/Library/Fonts/Hiragino Sans GB.ttc"]
    else:
        candidates = [os.path.expanduser("~/Library/Fonts/GoogleSans-Regular.ttf"), "/Library/Fonts/GoogleSans-Regular.ttf"]
    for p in candidates:
        if os.path.exists(p): return p
    fallback_all = ["/System/Library/Fonts/AppleSDGothicNeo.ttc", "/Library/Fonts/Arial Unicode.ttf"]
    for p in fallback_all:
        if os.path.exists(p): return p
    return None

def get_all_text_frames(shape):
    frames = []
    if getattr(shape, "has_text_frame", False): frames.append(shape.text_frame)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if getattr(cell, "text_frame", None): frames.append(cell.text_frame)
    if hasattr(shape, "shapes"):
        for child in shape.shapes: frames.extend(get_all_text_frames(child))
    return frames

def translate_batch_gemini(client, batch_texts, system_instruction, src_lang_name, lang_name):
    if not batch_texts: return batch_texts
    input_dict = {str(i): text for i, text in enumerate(batch_texts)}
    prompt = f"Translate this JSON object from {src_lang_name} into {lang_name}. Output ONLY a valid JSON object.\n\n" + json.dumps(input_dict, ensure_ascii=False)
    try:
        response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt, config=types.GenerateContentConfig(system_instruction=system_instruction, temperature=0.0, response_mime_type="application/json"))
        rj = json.loads(response.text)
        return [rj.get(str(i), batch_texts[i]) for i in range(len(batch_texts))]
    except Exception as e:
        if "spending cap" in str(e).lower() or "RESOURCE_EXHAUSTED" in str(e): raise SpendingCapExceeded(str(e))
        return batch_texts

def translate_batch_claude(client, batch_texts, system_instruction, src_lang_name="Korean", lang_name="Business English"):
    if not batch_texts: return batch_texts
    input_dict = {str(i): text for i, text in enumerate(batch_texts)}
    prompt = f"Translate this JSON object from {src_lang_name} into {lang_name}. Output ONLY a valid JSON object.\n\n" + json.dumps(input_dict, ensure_ascii=False)
    try:
        response = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=4096,
            temperature=0.0,
            system=system_instruction,
            messages=[{"role": "user", "content": prompt}]
        )
        response_text = response.content[0].text
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()
        rj = json.loads(response_text)
        return [rj.get(str(i), batch_texts[i]) for i in range(len(batch_texts))]
    except Exception as e:
        return batch_texts

class SpendingCapExceeded(Exception): pass

def translate_all_gemini_concurrent(client, texts, system_instruction, src_code, tgt_code, glossary, src_lang_name, lang_name, max_workers=5):
    total = len(texts)
    results = [None] * total
    batches = []
    for i in range(0, total, BATCH_SIZE):
        batches.append((i, texts[i:i+BATCH_SIZE]))
    
    spending_cap_hit = False
    
    def _worker(batch_info):
        nonlocal spending_cap_hit
        if spending_cap_hit: return None
        start_idx, batch_texts = batch_info
        try:
            res = translate_batch_gemini(client, batch_texts, system_instruction, src_lang_name, lang_name)
            if SLEEP_BETWEEN_BATCHES > 0: time.sleep(SLEEP_BETWEEN_BATCHES)
            return start_idx, res
        except SpendingCapExceeded:
            spending_cap_hit = True
            raise
        except Exception:
            return start_idx, batch_texts

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(_worker, b): b for b in batches}
        with tqdm(total=total, unit="항목", colour="green", desc="Gemini 번역 중") as pbar:
            for future in as_completed(futures):
                try:
                    result = future.result()
                    if result:
                        start_idx, translated_batch = result
                        for j, t in enumerate(translated_batch):
                            results[start_idx + j] = t
                        pbar.update(len(translated_batch))
                except SpendingCapExceeded:
                    # Cancel pending futures is not directly possible, but _worker checks spending_cap_hit
                    pass
    
    # Fill missing with free translator if Gemini failed or hit cap
    missing_indices = [i for i, r in enumerate(results) if r is None]
    if missing_indices:
        missing_texts = [texts[i] for i in missing_indices]
        logging.info(f"Gemini 번역 실패/중단으로 인해 {len(missing_texts)}개 항목을 무료 번역기로 전환합니다.")
        free_results = translate_all_free_concurrent(missing_texts, src_code, tgt_code, glossary)
        for idx, res in zip(missing_indices, free_results):
            results[idx] = res
            
    return results

def translate_all_claude_concurrent(client, texts, system_instruction, src_code, tgt_code, glossary, src_lang_name, lang_name, max_workers=5):
    total = len(texts)
    results = [None] * total
    batches = [(i, texts[i:i+BATCH_SIZE]) for i in range(0, total, BATCH_SIZE)]
    
    def _worker(batch_info):
        start_idx, batch_texts = batch_info
        try:
            res = translate_batch_claude(client, batch_texts, system_instruction, src_lang_name, lang_name)
            if SLEEP_BETWEEN_BATCHES > 0: time.sleep(SLEEP_BETWEEN_BATCHES)
            return start_idx, res
        except Exception:
            return start_idx, batch_texts

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(_worker, b): b for b in batches}
        with tqdm(total=total, unit="항목", colour="magenta", desc="Claude 번역 중") as pbar:
            for future in as_completed(futures):
                try:
                    start_idx, translated_batch = future.result()
                    for j, t in enumerate(translated_batch):
                        results[start_idx + j] = t
                    pbar.update(len(translated_batch))
                except: pass
    
    missing_indices = [i for i, r in enumerate(results) if r is None]
    if missing_indices:
        missing_texts = [texts[i] for i in missing_indices]
        free_results = translate_all_free_concurrent(missing_texts, src_code, tgt_code, glossary)
        for idx, res in zip(missing_indices, free_results):
            results[idx] = res
            
    return results

def translate_all_free_concurrent(texts, src_code, tgt_code, glossary, max_workers=20):
    def _one(args):
        idx, text = args
        temp = text
        for k, v in glossary.items(): temp = temp.replace(k, f"[{v}]")
        try:
            t = GoogleTranslator(source=src_code, target=tgt_code)
            res = t.translate(temp)
            return idx, res.replace("[", "").replace("]", "") if res else text
        except: return idx, text
    results = list(texts)
    n = len(texts)
    with ThreadPoolExecutor(max_workers=min(max_workers, n) if n > 0 else 1) as ex:
        futures = {ex.submit(_one, (i, t)): i for i, t in enumerate(texts)}
        with tqdm(total=n, unit="항목", colour="cyan", leave=False) as pbar:
            for f in as_completed(futures):
                idx, trans = f.result()
                results[idx] = trans
                pbar.update(1)
    return results

def shrink_text_gemini(client, text, target_len, lang_name):
    prompt = f"Compress this {lang_name} text strictly within {target_len} characters. Core meaning only. No explanation.\n\nText: {text}"
    try:
        response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt, config=types.GenerateContentConfig(temperature=0.0))
        return response.text.strip()
    except: return text

def shrink_text_claude(client, text, target_len, lang_name):
    prompt = f"Compress this {lang_name} text strictly within {target_len} characters. Core meaning only. No explanation.\n\nText: {text}"
    try:
        response = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=1024,
            temperature=0.0,
            messages=[{"role": "user", "content": prompt}]
        )
        return response.content[0].text.strip()
    except: return text

def shrink_all_gemini_concurrent(client, texts_to_shrink, target_lens, lang_name, max_workers=10):
    if not texts_to_shrink: return []
    results = list(texts_to_shrink)
    n = len(texts_to_shrink)
    
    def _worker(idx):
        return idx, shrink_text_gemini(client, texts_to_shrink[idx], target_lens[idx], lang_name)
    
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(_worker, i): i for i in range(n)}
        with tqdm(total=n, unit="항목", colour="yellow", desc="텍스트 압축 중", leave=False) as pbar:
            for future in as_completed(futures):
                idx, res = future.result()
                results[idx] = res
                pbar.update(1)
    return results

def shrink_all_claude_concurrent(client, texts_to_shrink, target_lens, lang_name, max_workers=10):
    if not texts_to_shrink: return []
    results = list(texts_to_shrink)
    n = len(texts_to_shrink)
    
    def _worker(idx):
        return idx, shrink_text_claude(client, texts_to_shrink[idx], target_lens[idx], lang_name)
    
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(_worker, i): i for i in range(n)}
        with tqdm(total=n, unit="항목", colour="yellow", desc="텍스트 압축 중(Claude)", leave=False) as pbar:
            for future in as_completed(futures):
                idx, res = future.result()
                results[idx] = res
                pbar.update(1)
    return results

def translate_pptx(input_path, output_path, client, system_instruction, src_code, tgt_code, src_regex, glossary, llm_provider, src_lang_name, lang_name, max_pages=None):
    prs = Presentation(input_path)
    paragraphs = []
    slides = list(prs.slides)
    if max_pages: slides = slides[:max_pages]
    for slide in slides:
        for shape in slide.shapes:
            for tf in get_all_text_frames(shape):
                for p in tf.paragraphs:
                    txt = p.text.strip()
                    if txt and re.search(src_regex, txt): paragraphs.append((p, txt))
    total = len(paragraphs)
    if total == 0: return 0
    else:
        if llm_provider == "Claude":
            translated = translate_all_claude_concurrent(client, [info[1] for info in paragraphs], system_instruction, src_code, tgt_code, glossary, src_lang_name, lang_name)
        elif llm_provider == "Gemini":
            translated = translate_all_gemini_concurrent(client, [info[1] for info in paragraphs], system_instruction, src_code, tgt_code, glossary, src_lang_name, lang_name)
        else:
            translated = translate_all_free_concurrent([info[1] for info in paragraphs], src_code, tgt_code, glossary)
    
    for (p, _), nt in zip(paragraphs, translated):
        if p.runs:
            p.runs[0].text = nt
            for j in range(len(p.runs)-1, 0, -1): p.runs[j]._r.getparent().remove(p.runs[j]._r)
        else: p.text = nt
    prs.save(output_path)
    return total

def translate_pdf(input_path, output_path, client, system_instruction, src_code, tgt_code, src_regex, glossary, llm_provider, src_lang_name, lang_name, max_pages=None):
    doc = fitz.open(input_path)
    font_path = find_google_sans_font(tgt_code)
    FONT_NAME = "f0"
    limit = max_pages if max_pages else len(doc)
    all_blocks = []
    for page_num in range(min(limit, len(doc))):
        page = doc[page_num]
        links = page.get_links()
        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b.get("type") != 0: continue
            current_p = {"text": "", "bbox": None, "fontsize": None, "color": None, "is_link": False}
            for line in b["lines"]:
                for span in line["spans"]:
                    stxt = span["text"].strip()
                    if not stxt: continue
                    srct = fitz.Rect(span["bbox"])
                    lfs = span["size"]
                    lclr = span["color"]
                    isl = any(fitz.Rect(l["from"]).intersects(srct) for l in links if l.get("kind") in (1, 2, 4))
                    isli = bool(re.match(r"^\s*(?:Step\s*\d+:?|\d+[\.\)]|\-|•|\*|[A-Za-z]\.|[①-⑳]|[⑴-⑽])\s+", stxt))
                    if current_p["fontsize"] is None:
                        current_p = {"text": span["text"], "bbox": srct, "fontsize": lfs, "color": lclr, "is_link": isl}
                    else:
                        color_diff = current_p["color"] != lclr
                        fs_diff = abs(current_p["fontsize"] - lfs) > 1.5
                        link_diff = current_p["is_link"] != isl
                        hg = srct.x0 - current_p["bbox"].x1
                        vg = srct.y0 - current_p["bbox"].y1
                        is_far_h = (hg > lfs * 1.5) and (abs(srct.y0 - current_p["bbox"].y0) < 5)
                        ends_dot = current_p["text"].strip().endswith(('。', '.', '！', '？', ':'))
                        is_far_v = (vg > lfs * 0.8) or (vg > 0 and ends_dot)
                        if (color_diff or fs_diff or link_diff or isli or is_far_h or is_far_v):
                            if current_p["text"].strip() and re.search(src_regex, current_p["text"]):
                                all_blocks.append({"page": page_num, "bbox": current_p["bbox"], "text": current_p["text"].strip(), "fontsize": current_p["fontsize"], "color": current_p["color"]})
                            current_p = {"text": span["text"], "bbox": srct, "fontsize": lfs, "color": lclr, "is_link": isl}
                        else:
                            ct = current_p["text"]
                            current_p["text"] = ct[:-1] + span["text"] if ct.endswith("-") else (ct + (" " if ct and not ct.endswith(" ") and not span["text"].startswith(" ") else "") + span["text"])
                            current_p["bbox"] |= srct
            if current_p["fontsize"] is not None and current_p["text"].strip() and re.search(src_regex, current_p["text"]):
                all_blocks.append({"page": page_num, "bbox": current_p["bbox"], "text": current_p["text"].strip(), "fontsize": current_p["fontsize"], "color": current_p["color"]})
    total = len(all_blocks)
    if total == 0: return 0
    if llm_provider == "Free (Google Translator)":
        trans_res = translate_all_free_concurrent([b["text"] for b in all_blocks], src_code, tgt_code, glossary)
    elif llm_provider == "Claude":
        trans_res = translate_all_claude_concurrent(client, [b["text"] for b in all_blocks], system_instruction, src_code, tgt_code, glossary, src_lang_name, lang_name)
    else:
        trans_res = translate_all_gemini_concurrent(client, [b["text"] for b in all_blocks], system_instruction, src_code, tgt_code, glossary, src_lang_name, lang_name)
    
    if llm_provider != "Free (Google Translator)":
        to_shrink_indices = []
        for i, (b, t) in enumerate(zip(all_blocks, trans_res)):
            if len(t) > len(b["text"]) * SHRINK_THRESHOLD:
                to_shrink_indices.append(i)
        
        if to_shrink_indices:
            texts_to_shrink = [trans_res[i] for i in to_shrink_indices]
            target_lens = [int(len(all_blocks[i]["text"]) * 1.15) for i in to_shrink_indices]
            if llm_provider == "Claude":
                shrunk_texts = shrink_all_claude_concurrent(client, texts_to_shrink, target_lens, lang_name)
            else:
                shrunk_texts = shrink_all_gemini_concurrent(client, texts_to_shrink, target_lens, lang_name)
            for i, shrunk in zip(to_shrink_indices, shrunk_texts):
                trans_res[i] = shrunk
    pg = defaultdict(list)
    for b, t in zip(all_blocks, trans_res): pg[b["page"]].append((b, t))
    for p_num, items in pg.items():
        page = doc[p_num]
        for b, _ in items: page.add_redact_annot(b["bbox"], fill=None)
        page.apply_redactions(images=0, graphics=0, text=0)
        c_font = "helv"
        if font_path:
            try: page.insert_font(fontname=FONT_NAME, fontfile=font_path); c_font = FONT_NAME
            except: pass
        for b, t in items:
            c = b["color"]; color = ((c>>16&0xFF)/255, (c>>8&0xFF)/255, (c&0xFF)/255)
            try:
                rect = b["bbox"]; ofs = b["fontsize"]; inserted = False
                for fs in range(int(math.ceil(ofs)), 4, -1):
                    shape = page.new_shape()
                    if shape.insert_textbox(rect, t, fontsize=fs, fontname=c_font, color=color, align=0) >= 0:
                        shape.commit(); inserted = True; break
                if not inserted:
                    shape = page.new_shape(); shape.insert_textbox(rect, t, fontsize=5, fontname=c_font, color=color, align=0); shape.commit()
            except: pass
    rt = output_path + ".render_tmp.pdf"
    doc.save(rt, garbage=4, deflate=True)
    doc.close()
    if not stitch_links_with_pikepdf(input_path, rt, output_path): os.rename(rt, output_path)
    else:
        if os.path.exists(rt): os.remove(rt)
    return total

def main():
    parser = argparse.ArgumentParser(description="Gartner AI PPT/PDF 번역기 (CLI)")
    parser.add_argument("input", help="파일 경로")
    parser.add_argument("--direction", default="한국어 → English", choices=list(DIRECTION_MAP.keys()))
    parser.add_argument("--provider", default="free", choices=["free", "gemini", "claude"])
    parser.add_argument("--test", action="store_true")
    parser.add_argument("--output", default=None)
    args = parser.parse_args()
    ip = Path(args.input)
    if not ip.exists(): print(f"❌ 파일 없음: {ip}"); sys.exit(1)
    load_dotenv()
    
    client = None
    llm_provider = "Free (Google Translator)"
    
    if args.provider == "gemini":
        ak = os.getenv("GEMINI_API_KEY")
        if not ak: print("⚠️ Gemini API Key 없음. --provider free 모드 권장."); sys.exit(1)
        client = genai.Client(api_key=ak)
        llm_provider = "Gemini"
    elif args.provider == "claude":
        ak = os.getenv("ANTHROPIC_API_KEY")
        if not ak: print("⚠️ Anthropic API Key 없음. --provider free 모드 권장."); sys.exit(1)
        client = anthropic.Anthropic(api_key=ak)
        llm_provider = "Claude"
    d = DIRECTION_MAP[args.direction]
    sc, tc, srcn, tgn, sre = d["src_code"], d["tgt_code"], d["src_lang_name"], d["lang_name"], d["src_regex"]
    gl = load_glossary()
    sys_ins = f"You are a Senior Gartner Analyst. Translate {srcn} to {tgn}. Concise bullet-points. Glossary: {json.dumps(gl)}"
    if args.output:
        op = args.output
    else:
        now_str = datetime.now().strftime("%Y%m%d_%H%M%S")
        result_dir = Path(__file__).resolve().parent.parent / "result"
        result_dir.mkdir(parents=True, exist_ok=True)
        op = str(result_dir / f"{ip.stem}_Translated_{now_str}{ip.suffix}")
    
    max_p = 20 if args.test else None
    if ip.suffix.lower() == ".pdf": translate_pdf(str(ip), op, client, sys_ins, sc, tc, sre, gl, llm_provider, srcn, tgn, max_p)
    else: translate_pptx(str(ip), op, client, sys_ins, sc, tc, sre, gl, llm_provider, srcn, tgn, max_p)
    print(f"✅ 완료: {op}")

if __name__ == "__main__": main()
