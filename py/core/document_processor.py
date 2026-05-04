import os
import re
import logging
from abc import ABC, abstractmethod

class DocumentProcessor(ABC):
    def __init__(self, translator):
        self.translator = translator

    @abstractmethod
    def process(self, input_path, output_path, progress_callback=None):
        pass

class PDFProcessor(DocumentProcessor):
    def process(self, input_path, output_path, progress_callback=None):
        import fitz
        import pikepdf
        
        logging.info(f"Starting PDF translation: {input_path}")
        try:
            doc = fitz.open(input_path)
            total_pages = len(doc)
            
            # Step 1: Extract Text
            if progress_callback: progress_callback(0, total_pages * 2, "텍스트 추출 중...")
            all_blocks = []
            for page_num in range(total_pages):
                page = doc[page_num]
                blocks = page.get_text("dict")["blocks"]
                for b in blocks:
                    if b["type"] == 0:
                        text = " ".join([l["spans"][0]["text"] for l in b["lines"]]).strip()
                        if text and re.search(self.translator.src_regex, text):
                            all_blocks.append({
                                "page": page_num,
                                "bbox": fitz.Rect(b["bbox"]),
                                "text": text
                            })
                            
            if not all_blocks:
                logging.info("번역할 텍스트가 없습니다.")
                doc.save(output_path)
                return True

            # Step 2: Translate
            if progress_callback: progress_callback(total_pages, total_pages * 2, "AI 번역 진행 중...")
            texts_to_translate = [b["text"] for b in all_blocks]
            
            # Use the translator's concurrent translation
            translated_texts = self.translator.translate_all_concurrent(texts_to_translate)
            
            # Replace empty translations with fallback text
            translated_texts = [t if t else "번역 실패" for t in translated_texts]
            
            for i, b in enumerate(all_blocks):
                b["translated"] = translated_texts[i]

            # Step 3: Draw Text
            if progress_callback: progress_callback(total_pages + (total_pages//2), total_pages * 2, "PDF 렌더링 중...")
            
            # Determine Font
            def find_font(tgt_code):
                base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
                google_sans = os.path.join(base_dir, "font", "Google Sans", "GoogleSans-Regular.ttf")
                if os.path.exists(google_sans): return google_sans
                if tgt_code == "ko":
                    candidates = ["/System/Library/Fonts/AppleSDGothicNeo.ttc", "/Library/Fonts/NanumGothic.ttf"]
                elif tgt_code == "ja":
                    candidates = ["/System/Library/Fonts/Hiragino Sans GB.ttc"]
                else:
                    candidates = ["/Library/Fonts/Arial Unicode.ttf"]
                for c in candidates:
                    if os.path.exists(c): return c
                return None

            font_path = find_font(self.translator.tgt_lang_code)
            
            for page_num in range(total_pages):
                page = doc[page_num]
                if font_path:
                    try:
                        page.insert_font(fontname="F0", fontfile=font_path)
                        fontname = "F0"
                    except:
                        fontname = "helv"
                else:
                    fontname = "helv"

                page_blocks = [b for b in all_blocks if b["page"] == page_num]
                for b in page_blocks:
                    rect = b["bbox"]
                    page.add_redact_annot(rect, fill=(1,1,1))
                    page.apply_redactions()
                    
                    text_to_insert = b["translated"]
                    try:
                        rc = page.insert_textbox(rect, text_to_insert, fontsize=11, fontname=fontname, align=0, color=(0,0,0))
                        if rc < 0:
                            shrunk = self.translator.shrink_text(text_to_insert, 100)
                            page.insert_textbox(rect, shrunk, fontsize=9, fontname=fontname, align=0, color=(0,0,0))
                    except Exception as e:
                        logging.warning(f"Text box insert failed: {e}")

            temp_path = output_path + ".tmp.pdf"
            doc.save(temp_path)
            doc.close()

            # Step 4: Stitch Links using pikepdf
            if progress_callback: progress_callback(total_pages * 2, total_pages * 2, "링크 복원 및 저장 중...")
            orig_pdf = pikepdf.Pdf.open(input_path)
            trans_pdf = pikepdf.Pdf.open(temp_path)
            
            for i, orig_page in enumerate(orig_pdf.pages):
                if i < len(trans_pdf.pages):
                    trans_page = trans_pdf.pages[i]
                    if "/Annots" in orig_page:
                        link_annots = pikepdf.Array()
                        for annot in orig_page.Annots:
                            if annot.get("/Subtype") == "/Link":
                                link_annots.append(trans_pdf.copy_foreign(annot))
                        if len(link_annots) > 0:
                            if "/Annots" not in trans_page:
                                trans_page.Annots = pikepdf.Array()
                            trans_page.Annots.extend(link_annots)
            
            # Fix global names (bookmarks)
            if "/Names" in orig_pdf.Root:
                trans_pdf.Root.Names = trans_pdf.copy_foreign(orig_pdf.Root.Names)

            trans_pdf.save(output_path)
            os.remove(temp_path)
            logging.info(f"PDF Translation complete: {output_path}")
            return True
            
        except Exception as e:
            logging.error(f"PDF processing failed: {e}")
            return False


class PPTXProcessor(DocumentProcessor):
    def process(self, input_path, output_path, progress_callback=None):
        from pptx import Presentation
        
        logging.info(f"Starting PPTX translation: {input_path}")
        try:
            prs = Presentation(input_path)
            total_slides = len(prs.slides)
            
            if progress_callback: progress_callback(0, total_slides * 2, "텍스트 추출 중...")
            
            def get_all_text_frames(shape):
                frames = []
                if getattr(shape, "has_text_frame", False): frames.append(shape.text_frame)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if getattr(cell, "text_frame", None): frames.append(cell.text_frame)
                if hasattr(shape, "shapes"):
                    for child in shape.shapes: frames.extend(get_all_text_frames(child))
                return frames

            all_paragraphs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    frames = get_all_text_frames(shape)
                    for frame in frames:
                        for p in frame.paragraphs:
                            text = p.text.strip()
                            if text and re.search(self.translator.src_regex, text):
                                all_paragraphs.append(p)

            if not all_paragraphs:
                logging.info("번역할 텍스트가 없습니다.")
                prs.save(output_path)
                return True

            if progress_callback: progress_callback(total_slides, total_slides * 2, "AI 번역 진행 중...")
            
            texts_to_translate = [p.text for p in all_paragraphs]
            translated_texts = self.translator.translate_all_concurrent(texts_to_translate)
            
            if progress_callback: progress_callback(total_slides + (total_slides//2), total_slides * 2, "PPT 렌더링 중...")

            for i, p in enumerate(all_paragraphs):
                t = translated_texts[i]
                if t:
                    if len(p.runs) > 0:
                        p.runs[0].text = t
                        for r in p.runs[1:]: r.text = ""
                    else:
                        p.text = t

            prs.save(output_path)
            if progress_callback: progress_callback(total_slides * 2, total_slides * 2, "완료")
            logging.info(f"PPTX Translation complete: {output_path}")
            return True
            
        except Exception as e:
            logging.error(f"PPTX processing failed: {e}")
            return False
